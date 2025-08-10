"""
Document Analysis Plugin
Handles multi-format document parsing, summarization, and Q&A
"""

import os
import asyncio
import base64
from typing import Dict, List, Any, Optional
from pathlib import Path
import tempfile

# Document processing libraries
import PyPDF2
from docx import Document
import pandas as pd
from pptx import Presentation
from PIL import Image
import io

# Local AI processing
from transformers import pipeline, AutoTokenizer, AutoModelForQuestionAnswering
import torch

from . import PluginBase

class DocumentAnalysisPlugin(PluginBase):
    """Plugin for document analysis and processing"""
    
    def __init__(self):
        self.name = "document_analysis"
        self.supported_formats = {
            '.pdf': self._parse_pdf,
            '.docx': self._parse_docx,
            '.txt': self._parse_txt,
            '.md': self._parse_txt,
            '.rtf': self._parse_txt,
            '.csv': self._parse_csv,
            '.xlsx': self._parse_excel,
            '.pptx': self._parse_pptx,
            '.png': self._parse_image,
            '.jpg': self._parse_image,
            '.jpeg': self._parse_image
        }
        
        # Initialize local AI models
        self.summarizer = None
        self.qa_model = None
        self.tokenizer = None
        self._initialize_models()
        
        # Document storage
        self.documents = {}
    
    def _initialize_models(self):
        """Initialize local AI models for summarization and Q&A"""
        try:
            # Use smaller models for local processing
            self.summarizer = pipeline(
                "summarization",
                model="facebook/bart-large-cnn",
                device=0 if torch.cuda.is_available() else -1
            )
            
            model_name = "deepset/roberta-base-squad2"
            self.tokenizer = AutoTokenizer.from_pretrained(model_name)
            self.qa_model = AutoModelForQuestionAnswering.from_pretrained(model_name)
            
            if torch.cuda.is_available():
                self.qa_model = self.qa_model.cuda()
                
            print("Document analysis models loaded successfully")
        except Exception as e:
            print(f"Warning: Could not load AI models: {e}")
    
    def get_name(self) -> str:
        return self.name
    
    def get_description(self) -> str:
        return "Document analysis and processing with multi-format support"
    
    def get_functions(self) -> List[Dict[str, Any]]:
        return [
            {
                "name": "analyze_document",
                "description": "Analyze and extract content from uploaded documents (PDF, DOCX, TXT, MD, RTF, CSV, XLSX, PPTX, PNG, JPG)",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "file_data": {"type": "string", "description": "Base64 encoded file data"},
                        "filename": {"type": "string", "description": "Original filename with extension"},
                        "analysis_type": {"type": "string", "description": "Type of analysis: 'extract', 'summarize', or 'qa'"},
                        "question": {"type": "string", "description": "Question for Q&A analysis (required if analysis_type is 'qa')"}
                    },
                    "required": ["file_data", "filename", "analysis_type"]
                }
            },
            {
                "name": "list_documents",
                "description": "List all uploaded and processed documents",
                "parameters": {
                    "type": "object",
                    "properties": {},
                    "required": []
                }
            },
            {
                "name": "delete_document",
                "description": "Delete a document from storage",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "document_id": {"type": "string", "description": "ID of document to delete"}
                    },
                    "required": ["document_id"]
                }
            }
        ]
    
    async def execute_function(self, function_name: str, **kwargs) -> Dict[str, Any]:
        if function_name == "analyze_document":
            return await self._analyze_document(**kwargs)
        elif function_name == "list_documents":
            return await self._list_documents(**kwargs)
        elif function_name == "delete_document":
            return await self._delete_document(**kwargs)
        else:
            return {"error": f"Unknown function: {function_name}"}
    
    def get_widget_info(self) -> Optional[Dict[str, Any]]:
        return {
            "component_name": "DocumentAnalysisWidget",
            "css_file": "DocumentAnalysisWidget.css",
            "socket_events": ["document_analysis_update", "document_list_update"]
        }
    
    async def _analyze_document(self, file_data: str, filename: str, analysis_type: str, question: str = None) -> Dict[str, Any]:
        """Analyze uploaded document"""
        try:
            # Decode base64 file data
            file_bytes = base64.b64decode(file_data)
            
            # Get file extension
            file_ext = Path(filename).suffix.lower()
            
            if file_ext not in self.supported_formats:
                return {"error": f"Unsupported file format: {file_ext}"}
            
            # Parse document content
            content = await self._parse_document(file_bytes, file_ext)
            
            if not content:
                return {"error": "Failed to extract content from document"}
            
            # Generate document ID
            import hashlib
            doc_id = hashlib.md5(file_bytes).hexdigest()
            
            # Store document
            self.documents[doc_id] = {
                "filename": filename,
                "content": content,
                "file_size": len(file_bytes),
                "upload_time": asyncio.get_event_loop().time()
            }
            
            # Perform analysis
            result = {
                "document_id": doc_id,
                "filename": filename,
                "content_length": len(content),
                "analysis_type": analysis_type
            }
            
            if analysis_type == "extract":
                result["extracted_content"] = content[:2000] + "..." if len(content) > 2000 else content
                
            elif analysis_type == "summarize":
                if self.summarizer:
                    summary = await self._summarize_text(content)
                    result["summary"] = summary
                else:
                    result["summary"] = "Summarization not available (AI models not loaded)"
                    
            elif analysis_type == "qa" and question:
                if self.qa_model and self.tokenizer:
                    answer = await self._answer_question(content, question)
                    result["question"] = question
                    result["answer"] = answer
                else:
                    result["answer"] = "Q&A not available (AI models not loaded)"
            
            return result
            
        except Exception as e:
            return {"error": f"Document analysis failed: {str(e)}"}
    
    async def _parse_document(self, file_bytes: bytes, file_ext: str) -> str:
        """Parse document based on file extension"""
        try:
            # Create temporary file
            with tempfile.NamedTemporaryFile(suffix=file_ext, delete=False) as temp_file:
                temp_file.write(file_bytes)
                temp_file_path = temp_file.name
            
            # Parse using appropriate method
            parser_func = self.supported_formats.get(file_ext)
            if parser_func:
                content = await asyncio.to_thread(parser_func, temp_file_path)
            else:
                content = "Unsupported file format"
            
            # Clean up temporary file
            os.unlink(temp_file_path)
            
            return content
            
        except Exception as e:
            print(f"Error parsing document: {e}")
            return ""
    
    def _parse_pdf(self, file_path: str) -> str:
        """Parse PDF file"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text
        except Exception as e:
            print(f"PDF parsing error: {e}")
            return ""
    
    def _parse_docx(self, file_path: str) -> str:
        """Parse DOCX file"""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            print(f"DOCX parsing error: {e}")
            return ""
    
    def _parse_txt(self, file_path: str) -> str:
        """Parse text files (TXT, MD, RTF)"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                print(f"Text file parsing error: {e}")
                return ""
    
    def _parse_csv(self, file_path: str) -> str:
        """Parse CSV file"""
        try:
            df = pd.read_csv(file_path)
            return df.to_string()
        except Exception as e:
            print(f"CSV parsing error: {e}")
            return ""
    
    def _parse_excel(self, file_path: str) -> str:
        """Parse Excel file"""
        try:
            df = pd.read_excel(file_path)
            return df.to_string()
        except Exception as e:
            print(f"Excel parsing error: {e}")
            return ""
    
    def _parse_pptx(self, file_path: str) -> str:
        """Parse PowerPoint file"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"PowerPoint parsing error: {e}")
            return ""
    
    def _parse_image(self, file_path: str) -> str:
        """Parse image file (OCR would be added here)"""
        try:
            # For now, return basic image info
            with Image.open(file_path) as img:
                return f"Image: {img.format} {img.size} {img.mode}"
        except Exception as e:
            print(f"Image parsing error: {e}")
            return ""
    
    async def _summarize_text(self, text: str) -> str:
        """Summarize text using local AI model"""
        try:
            # Truncate text if too long
            max_length = 1024
            if len(text) > max_length:
                text = text[:max_length]
            
            summary = await asyncio.to_thread(
                self.summarizer,
                text,
                max_length=150,
                min_length=50,
                do_sample=False
            )
            
            return summary[0]['summary_text']
        except Exception as e:
            print(f"Summarization error: {e}")
            return "Summarization failed"
    
    async def _answer_question(self, context: str, question: str) -> str:
        """Answer question using local AI model"""
        try:
            # Truncate context if too long
            max_length = 512
            if len(context) > max_length:
                context = context[:max_length]
            
            inputs = await asyncio.to_thread(
                self.tokenizer,
                question,
                context,
                return_tensors="pt",
                max_length=512,
                truncation=True
            )
            
            if torch.cuda.is_available():
                inputs = {k: v.cuda() for k, v in inputs.items()}
            
            with torch.no_grad():
                outputs = self.qa_model(**inputs)
            
            answer_start = torch.argmax(outputs.start_logits)
            answer_end = torch.argmax(outputs.end_logits) + 1
            
            answer = self.tokenizer.convert_tokens_to_string(
                self.tokenizer.convert_ids_to_tokens(
                    inputs["input_ids"][0][answer_start:answer_end]
                )
            )
            
            return answer if answer else "No answer found"
            
        except Exception as e:
            print(f"Q&A error: {e}")
            return "Q&A failed"
    
    async def _list_documents(self) -> Dict[str, Any]:
        """List all stored documents"""
        docs = []
        for doc_id, doc_info in self.documents.items():
            docs.append({
                "id": doc_id,
                "filename": doc_info["filename"],
                "file_size": doc_info["file_size"],
                "upload_time": doc_info["upload_time"]
            })
        
        return {"documents": docs}
    
    async def _delete_document(self, document_id: str) -> Dict[str, Any]:
        """Delete a document from storage"""
        if document_id in self.documents:
            del self.documents[document_id]
            return {"success": True, "message": "Document deleted"}
        else:
            return {"error": "Document not found"}
